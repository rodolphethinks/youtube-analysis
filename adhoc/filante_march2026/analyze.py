"""
Ad-hoc Analysis: Renault Filante Influencer Campaign - March 2026
-----------------------------------------------------------------
Reads YouTube video URLs from the campaign Excel file, fetches all comments
via the YouTube Data API, classifies each comment sentiment with Gemini AI,
generates key takeaways, and produces:
  - output/comments_by_likes.csv   – all comments sorted by likes desc
  - output/report.html             – slide-style HTML dashboard
"""

import os
import sys
import re
import json
import time
from pathlib import Path

import pandas as pd
import openpyxl
from dotenv import load_dotenv

# ── project root on path ────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent.parent
sys.path.insert(0, str(ROOT))
load_dotenv(ROOT / ".env")

from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import google.generativeai as genai

# ── config ───────────────────────────────────────────────────────────────────
EXCEL_PATH = ROOT / "웰컴_RENAULT FILANTE Influencer Test Drive Content Coverage_260316 1.xlsx"
OUTPUT_DIR = Path(__file__).parent / "output"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
if not GOOGLE_API_KEY:
    raise EnvironmentError("GOOGLE_API_KEY is not set. Add it to your .env file.")

MAX_COMMENTS_PER_VIDEO = 100   # YouTube API page = max 100
GEMINI_MODEL = "gemini-2.0-flash"
SENTIMENT_BATCH_SIZE = 40      # comments per Gemini call


# ── Step 1: Extract YouTube video URLs from Excel ────────────────────────────

def clean_youtube_url(raw_url: str) -> str | None:
    """Normalise any YouTube URL form to https://www.youtube.com/watch?v=ID."""
    if not raw_url:
        return None
    # Skip non-YouTube URLs (blogs, Naver, etc.)
    if "youtube.com" not in raw_url and "youtu.be" not in raw_url:
        return None

    # Extract video ID – covers watch, shorts, youtu.be
    patterns = [
        r"youtube\.com/watch\?v=([A-Za-z0-9_-]{11})",
        r"youtube\.com/shorts/([A-Za-z0-9_-]{11})",
        r"youtu\.be/([A-Za-z0-9_-]{11})",
    ]
    for pat in patterns:
        m = re.search(pat, raw_url)
        if m:
            return f"https://www.youtube.com/watch?v={m.group(1)}"
    return None


def extract_videos_from_excel(excel_path: Path) -> list[dict]:
    """
    Parse the first sheet and collect one record per YouTube video.
    Returns list of {title, channel, url, video_id, views, likes, comments}.
    """
    wb = openpyxl.load_workbook(str(excel_path))
    ws = wb["콘텐츠 업로드 현황"]

    seen_ids: set[str] = set()
    videos: list[dict] = []

    # Rows 4-53 contain video data (1-indexed); header on row 2
    for row in ws.iter_rows(min_row=4, max_row=53):
        title_cell = row[4]  # column E (index 4)
        raw_url = title_cell.hyperlink.target if title_cell.hyperlink else None
        url = clean_youtube_url(raw_url or "")
        if not url:
            continue  # skip blogs and un-linked rows

        video_id = re.search(r"v=([A-Za-z0-9_-]{11})", url)
        if not video_id:
            continue
        vid = video_id.group(1)
        if vid in seen_ids:
            continue
        seen_ids.add(vid)

        # Channel is in col D (index 3)
        channel = row[3].value or ""
        # View count col G (index 6); likes col H (index 7); comments col I (index 8)
        views = row[6].value if isinstance(row[6].value, int) else 0
        likes = row[7].value if isinstance(row[7].value, int) else 0
        n_comments = row[8].value if isinstance(row[8].value, int) else 0

        videos.append({
            "title": str(title_cell.value or "").strip(),
            "channel": str(channel).strip(),
            "url": url,
            "video_id": vid,
            "views": views,
            "likes": likes,
            "comment_count": n_comments,
        })

    print(f"[Excel] Extracted {len(videos)} unique YouTube videos.")
    return videos


# ── Step 2: Fetch comments via YouTube Data API ──────────────────────────────

def build_youtube_client(api_key: str):
    return build("youtube", "v3", developerKey=api_key)


def fetch_comments(yt, video_id: str, max_comments: int = MAX_COMMENTS_PER_VIDEO) -> list[dict]:
    """Fetch top-level comments for a video ordered by relevance."""
    comments: list[dict] = []
    try:
        request = yt.commentThreads().list(
            part="snippet",
            videoId=video_id,
            maxResults=min(max_comments, 100),
            order="relevance",
            textFormat="plainText",
        )
        response = request.execute()
        for item in response.get("items", []):
            s = item["snippet"]["topLevelComment"]["snippet"]
            comments.append({
                "video_id": video_id,
                "author": s.get("authorDisplayName", ""),
                "comment": s.get("textDisplay", ""),
                "likes": int(s.get("likeCount", 0)),
                "published_at": s.get("publishedAt", "")[:10],
            })
    except HttpError as e:
        code = e.resp.status if hasattr(e, "resp") else "?"
        reason = str(e)
        if "commentsDisabled" in reason or code == 403:
            print(f"  [skip] Comments disabled for {video_id}")
        else:
            print(f"  [warn] YouTube API error for {video_id}: {e}")
    return comments


# ── Step 3: Gemini – batch comment sentiment classification ──────────────────

def init_gemini(api_key: str):
    genai.configure(api_key=api_key)
    return genai.GenerativeModel(GEMINI_MODEL)


def classify_sentiments_batch(model, comments: list[str]) -> list[str]:
    """
    Send a batch of comments to Gemini and return a list of sentiments
    ('Positive', 'Neutral', or 'Negative') in the same order.
    """
    numbered = "\n".join(f"{i+1}. {c}" for i, c in enumerate(comments))
    prompt = f"""You are a sentiment analysis expert specialised in automotive consumer opinions.

Classify each of the following YouTube comments about the Renault Filante car as exactly one of:
  Positive | Neutral | Negative

Rules:
- Base your classification only on the comment content.
- Treat factual questions or ambiguous comments as Neutral.
- Return ONLY a valid JSON array of strings, one per comment, in the same order.
  Example: ["Positive", "Neutral", "Negative"]

Comments:
{numbered}"""

    try:
        resp = model.generate_content(prompt)
        text = resp.text.strip()
        # Strip markdown code fences if present
        text = re.sub(r"^```[a-z]*\n?", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\n?```$", "", text)
        result = json.loads(text)
        if isinstance(result, list) and len(result) == len(comments):
            return [s if s in ("Positive", "Neutral", "Negative") else "Neutral" for s in result]
    except Exception as e:
        print(f"  [warn] Gemini batch error: {e}")
    # Fallback: neutral for all
    return ["Neutral"] * len(comments)


def classify_all_comments(model, comments_df: pd.DataFrame) -> pd.DataFrame:
    """Add a 'sentiment' column to the comments DataFrame."""
    texts = comments_df["comment"].tolist()
    sentiments: list[str] = []
    total = len(texts)
    for start in range(0, total, SENTIMENT_BATCH_SIZE):
        batch = texts[start: start + SENTIMENT_BATCH_SIZE]
        print(f"  Classifying comments {start+1}–{min(start+len(batch), total)} / {total}…")
        sentiments.extend(classify_sentiments_batch(model, batch))
        time.sleep(0.5)  # polite rate limiting
    comments_df = comments_df.copy()
    comments_df["sentiment"] = sentiments
    return comments_df


# ── Step 4: Gemini – key takeaways from all comments ─────────────────────────

def generate_key_takeaways(model, comments_df: pd.DataFrame, max_comments: int = 600) -> list[dict]:
    """
    Ask Gemini to distil 4–5 key takeaways from the full comment corpus.
    Returns list of {title, body}.
    """
    # Sample up to max_comments, weighted by likes (high-like comments first)
    sample = comments_df.sort_values("likes", ascending=False).head(max_comments)
    corpus = "\n".join(f"- {row['comment']}" for _, row in sample.iterrows())

    prompt = f"""You are an automotive market research analyst.

Below are YouTube comments collected from influencer test-drive videos for the Renault Filante hybrid.
Based on all the comments, write 4–5 concise key takeaways that a marketing team would find useful.
Each takeaway should have a short title (3–5 words) and a 1–2 sentence explanation.

Return ONLY valid JSON as a list of objects: [{{"title": "...", "body": "..."}}]

Comments:
{corpus}"""

    try:
        resp = model.generate_content(prompt)
        text = resp.text.strip()
        text = re.sub(r"^```[a-z]*\n?", "", text, flags=re.IGNORECASE)
        text = re.sub(r"\n?```$", "", text)
        result = json.loads(text)
        if isinstance(result, list):
            return result
    except Exception as e:
        print(f"  [warn] Gemini takeaways error: {e}")
    return [{"title": "Analysis complete", "body": "See comments CSV for full details."}]


# ── Step 5: Generate PPTX report ─────────────────────────────────────────────

def _pt(val: float) -> int:
    """Convert points to EMU (English Metric Units) for python-pptx."""
    from pptx.util import Pt
    return Pt(val)


def _make_pie_chart_image(pos: int, neu: int, neg: int, out_path: Path) -> Path:
    """Render a matplotlib pie chart and save as PNG for embedding."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    total = pos + neu + neg
    sizes = [pos, neu, neg]
    colors = ["#27ae60", "#bdc3c7", "#e74c3c"]
    labels = [
        f"Positive\n{round(pos/total*100)}%",
        f"Neutral\n{round(neu/total*100)}%",
        f"Negative\n{round(neg/total*100)}%",
    ]

    fig, ax = plt.subplots(figsize=(4, 4), dpi=150)
    wedges, texts = ax.pie(
        sizes,
        colors=colors,
        startangle=90,
        wedgeprops={"linewidth": 2, "edgecolor": "white"},
    )
    ax.legend(
        wedges,
        labels,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.22),
        ncol=3,
        fontsize=8,
        frameon=False,
    )
    plt.tight_layout()
    fig.savefig(str(out_path), transparent=True, bbox_inches="tight")
    plt.close(fig)
    return out_path


def build_pptx_report(
    videos: list[dict],
    comments_df: pd.DataFrame,
    takeaways: list[dict],
    output_path: Path,
):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    # ── Colour palette ────────────────────────────────────────────────────────
    YELLOW  = RGBColor(0xF5, 0xC4, 0x00)
    BLACK   = RGBColor(0x11, 0x11, 0x11)
    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    LGRAY   = RGBColor(0xF4, 0xF4, 0xF4)
    MGRAY   = RGBColor(0xCC, 0xCC, 0xCC)
    DGRAY   = RGBColor(0x55, 0x55, 0x55)
    GREEN   = RGBColor(0x27, 0xAE, 0x60)
    RED     = RGBColor(0xE7, 0x4C, 0x3C)

    # ── Helpers ───────────────────────────────────────────────────────────────
    def add_rect(slide, l, t, w, h, fill=None, line=None, line_w_pt=0.5):
        from pptx.util import Pt as _Pt
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.line.fill.background()  # no line by default
        if fill:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
        else:
            shape.fill.background()
        if line:
            shape.line.color.rgb = line
            shape.line.width = _Pt(line_w_pt)
        return shape

    def add_text(slide, text, l, t, w, h, size=11, bold=False, color=BLACK,
                 align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        return tb

    def section_title(slide, title, l, t, w):
        """Small section header with yellow underline."""
        add_text(slide, title, l, t, w, 0.22, size=9, bold=True, color=DGRAY)
        add_rect(slide, l, t + 0.20, w, 0.03, fill=YELLOW)

    # ── Stats ─────────────────────────────────────────────────────────────────
    total_videos   = len(videos)
    total_comments = len(comments_df)
    counts  = comments_df["sentiment"].value_counts()
    pos = int(counts.get("Positive", 0))
    neu = int(counts.get("Neutral",  0))
    neg = int(counts.get("Negative", 0))
    pos_pct = round(pos / total_comments * 100) if total_comments else 0

    # ── Presentation setup (16:9 widescreen) ──────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # completely blank
    slide = prs.slides.add_slide(blank_layout)

    W = 13.33  # slide width  (inches)
    H = 7.5   # slide height (inches)

    # ── Background ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, fill=LGRAY)

    # ── Header bar ────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, 0.90, fill=WHITE)
    add_rect(slide, 0, 0, W, 0.07, fill=YELLOW)  # top accent bar
    add_text(slide, "Social listening", 0.3, 0.10, 3, 0.22,
             size=8, bold=True, color=YELLOW)
    add_text(slide, "SOCIAL MENTION ANALYSIS", 0.3, 0.30, 6, 0.45,
             size=22, bold=True, color=BLACK)
    # Period badge (dark pill)
    add_rect(slide, W - 2.2, 0.22, 2.0, 0.40, fill=BLACK)
    add_text(slide, "YouTube  ·  Q1 2026",
             W - 2.15, 0.26, 1.9, 0.32,
             size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── KPI cards  (three equal cards below header) ───────────────────────────
    kpis = [
        (f"{total_videos:,}",   "Videos Analyzed",       BLACK),
        (f"{total_comments:,}", "Total Comments Analyzed", BLACK),
        (f"{pos_pct}%",         "Positive Sentiment",      GREEN),
    ]
    card_w = 2.8
    card_gap = 0.18
    cards_total = card_w * 3 + card_gap * 2
    card_x0 = (W - cards_total) / 2
    card_y = 1.02

    for i, (val, lbl, val_color) in enumerate(kpis):
        cx = card_x0 + i * (card_w + card_gap)
        add_rect(slide, cx, card_y, card_w, 0.88, fill=WHITE, line=MGRAY)
        add_rect(slide, cx, card_y, card_w, 0.06, fill=YELLOW)  # top accent
        add_text(slide, val, cx + 0.12, card_y + 0.10, card_w - 0.24, 0.42,
                 size=26, bold=True, color=val_color)
        add_text(slide, lbl, cx + 0.12, card_y + 0.56, card_w - 0.24, 0.28,
                 size=9, color=DGRAY)

    # ── Three-column body ─────────────────────────────────────────────────────
    body_y = 2.08
    body_h = H - body_y - 0.22

    col_widths = [2.55, 5.70, 2.85]
    gutter = 0.13
    col_x = [0.13]
    for cw in col_widths[:-1]:
        col_x.append(col_x[-1] + cw + gutter)

    # Draw white column panels
    for i, cw in enumerate(col_widths):
        add_rect(slide, col_x[i], body_y, cw, body_h, fill=WHITE, line=MGRAY)

    # ── LEFT column: Sentiment ────────────────────────────────────────────────
    lx, lw = col_x[0], col_widths[0]
    section_title(slide, "Overall Sentiment", lx + 0.12, body_y + 0.12, lw - 0.24)

    # Pie chart image
    pie_img = OUTPUT_DIR / "_pie_tmp.png"
    _make_pie_chart_image(pos, neu, neg, pie_img)
    pie_size = 2.10
    slide.shapes.add_picture(
        str(pie_img),
        Inches(lx + (lw - pie_size) / 2),
        Inches(body_y + 0.45),
        Inches(pie_size), Inches(pie_size)
    )

    # Big % box
    big_y = body_y + body_h - 0.88
    add_rect(slide, lx + 0.12, big_y, lw - 0.24, 0.76, fill=RGBColor(0xE8, 0xF8, 0xEF), line=MGRAY)
    add_text(slide, f"{pos_pct}%", lx + 0.12, big_y + 0.04, lw - 0.24, 0.42,
             size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "of comments are Positive",
             lx + 0.12, big_y + 0.46, lw - 0.24, 0.22,
             size=8, color=DGRAY, align=PP_ALIGN.CENTER)

    # ── MIDDLE column: Comment highlights ─────────────────────────────────────
    mx, mw = col_x[1], col_widths[1]
    section_title(slide, "Comment Highlights", mx + 0.12, body_y + 0.12, mw - 0.24)

    top_comments = comments_df.sort_values("likes", ascending=False).head(4)
    card_h  = (body_h - 0.42) / 4 - 0.06
    sent_colors = {"Positive": GREEN, "Neutral": DGRAY, "Negative": RED}
    sent_bg = {
        "Positive": RGBColor(0xE8, 0xF8, 0xEF),
        "Neutral":  RGBColor(0xF0, 0xF0, 0xF0),
        "Negative": RGBColor(0xFD, 0xEC, 0xEA),
    }

    for idx, (_, row) in enumerate(top_comments.iterrows()):
        cy = body_y + 0.42 + idx * (card_h + 0.06)
        add_rect(slide, mx + 0.10, cy, mw - 0.20, card_h, fill=LGRAY, line=MGRAY)

        # author + likes
        add_text(slide, row["author"],
                 mx + 0.18, cy + 0.04, mw * 0.55, 0.18,
                 size=8, bold=True, color=BLACK)
        add_text(slide, f"{row['published_at']}   👍 {row['likes']}",
                 mx + mw * 0.55, cy + 0.04, mw * 0.40, 0.18,
                 size=7, color=DGRAY, align=PP_ALIGN.RIGHT)

        # comment text
        snippet = str(row["comment"])[:160] + ("…" if len(str(row["comment"])) > 160 else "")
        add_text(slide, f'"{snippet}"',
                 mx + 0.18, cy + 0.20, mw - 0.36, card_h - 0.36,
                 size=8, color=DGRAY, wrap=True)

        # sentiment pill
        sentiment = row.get("sentiment", "Neutral")
        sc = sent_colors.get(sentiment, DGRAY)
        sb = sent_bg.get(sentiment, LGRAY)
        pill = add_rect(slide, mx + 0.18, cy + card_h - 0.26, 0.72, 0.20, fill=sb)
        add_text(slide, sentiment,
                 mx + 0.18, cy + card_h - 0.26, 0.72, 0.20,
                 size=7, bold=True, color=sc, align=PP_ALIGN.CENTER)

    # ── RIGHT column: Key Takeaways ────────────────────────────────────────────
    rx, rw = col_x[2], col_widths[2]
    section_title(slide, "Key Takeaways", rx + 0.12, body_y + 0.12, rw - 0.24)

    item_y = body_y + 0.46
    for t in takeaways[:5]:
        title = t.get("title", "")
        body  = t.get("body",  "")
        # yellow dot
        dot = slide.shapes.add_shape(9, Inches(rx + 0.15), Inches(item_y + 0.06),
                                     Inches(0.13), Inches(0.13))
        dot.fill.solid()
        dot.fill.fore_color.rgb = YELLOW
        dot.line.fill.background()
        add_text(slide, title, rx + 0.34, item_y, rw - 0.46, 0.20,
                 size=9, bold=True, color=BLACK)
        add_text(slide, body,  rx + 0.34, item_y + 0.20, rw - 0.46, 0.46,
                 size=8, color=DGRAY, wrap=True)
        item_y += 0.80
        if item_y > body_y + body_h - 0.55:
            break

    # ── Footer ────────────────────────────────────────────────────────────────
    add_text(
        slide,
        f"Generated {pd.Timestamp.now().strftime('%B %d, %Y')}  |  Renault Korea – Filante Influencer Campaign",
        0.13, H - 0.20, W - 0.26, 0.18,
        size=7, color=MGRAY, align=PP_ALIGN.RIGHT,
    )

    prs.save(str(output_path))
    print(f"[Report] PPTX report saved → {output_path}")

    # Clean up temp pie image
    if pie_img.exists():
        pie_img.unlink()


# ── Main ──────────────────────────────────────────────────────────────────────

def _deleted_html():
    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8"/>
<meta name="viewport" content="width=device-width, initial-scale=1.0"/>
<title>Renault Filante – Social Mention Analysis</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{ font-family: 'Segoe UI', Arial, sans-serif; background: #f4f4f4; color: #222; }}

  /* ── Header ── */
  .page-header {{
    display: flex; justify-content: space-between; align-items: center;
    background: #fff; border-top: 5px solid #f5c400; padding: 12px 32px;
  }}
  .brand-label {{ font-size: 11px; color: #f5c400; font-weight: 700; letter-spacing: 1px; text-transform: uppercase; }}
  .page-title {{ font-size: 26px; font-weight: 900; letter-spacing: 1px; text-transform: uppercase; margin-top: 2px; }}
  .period-badge {{ background: #111; color: #fff; font-size: 12px; font-weight: 700; padding: 6px 16px; border-radius: 3px; }}

  /* ── KPI strip ── */
  .kpi-row {{ display: flex; gap: 16px; padding: 24px 32px 0; }}
  .kpi-card {{
    flex: 1; background: #fff; border: 1px solid #e0e0e0;
    padding: 20px 24px; border-top: 4px solid #f5c400;
  }}
  .kpi-value {{ font-size: 36px; font-weight: 900; line-height: 1.1; }}
  .kpi-value.positive {{ color: #27ae60; }}
  .kpi-label {{ font-size: 12px; color: #888; margin-top: 4px; }}

  /* ── Main panels ── */
  .panels {{ display: flex; gap: 16px; padding: 20px 32px 32px; align-items: flex-start; }}

  /* Left: sentiment */
  .panel-sentiment {{ flex: 0 0 260px; background: #fff; border: 1px solid #e0e0e0; padding: 20px; }}
  .panel-title {{ font-size: 13px; font-weight: 700; margin-bottom: 16px; color: #333; }}
  .chart-wrap {{ width: 200px; margin: 0 auto; }}
  .legend {{ display: flex; gap: 12px; justify-content: center; margin-top: 10px; flex-wrap: wrap; }}
  .legend-item {{ display: flex; align-items: center; gap: 4px; font-size: 11px; }}
  .legend-dot {{ width: 10px; height: 10px; border-radius: 50%; }}
  .big-pct {{
    text-align: center; margin-top: 16px; background: #e8f8ef;
    padding: 14px; font-size: 32px; font-weight: 900; color: #27ae60;
  }}
  .big-pct-label {{ font-size: 11px; color: #555; margin-top: 2px; }}

  /* Middle: comment highlights */
  .panel-comments {{ flex: 1; background: #fff; border: 1px solid #e0e0e0; padding: 20px; }}
  .comment-card {{
    border: 1px solid #e8e8e8; border-radius: 4px; padding: 12px 14px; margin-bottom: 10px;
  }}
  .comment-header {{ display: flex; justify-content: space-between; margin-bottom: 2px; }}
  .comment-author {{ font-size: 12px; font-weight: 700; color: #444; }}
  .comment-meta {{ font-size: 11px; color: #999; }}
  .comment-video-title {{ font-size: 10px; color: #aaa; margin-bottom: 6px; }}
  .comment-text {{ font-size: 13px; color: #333; font-style: italic; line-height: 1.5; }}
  .sentiment-badge {{
    display: inline-block; font-size: 10px; font-weight: 700;
    padding: 2px 8px; border-radius: 12px; margin-top: 6px; text-transform: uppercase;
  }}
  .sentiment-positive {{ background: #e8f8ef; color: #27ae60; }}
  .sentiment-negative {{ background: #fdecea; color: #c0392b; }}
  .sentiment-neutral  {{ background: #f0f0f0; color: #777; }}

  /* Right: key takeaways */
  .panel-takeaways {{ flex: 0 0 280px; background: #fff; border: 1px solid #e0e0e0; padding: 20px; }}
  .takeaway-item {{ display: flex; gap: 10px; margin-bottom: 18px; }}
  .takeaway-dot {{ width: 12px; height: 12px; background: #f5c400; border-radius: 50%; flex-shrink: 0; margin-top: 3px; }}
  .takeaway-title {{ font-size: 13px; font-weight: 700; margin-bottom: 4px; }}
  .takeaway-body {{ font-size: 12px; color: #555; line-height: 1.5; }}

  /* Footer */
  .footer {{ text-align: right; padding: 0 32px 16px; font-size: 10px; color: #aaa; }}
</style>
</head>
<body>

<div class="page-header">
  <div>
    <div class="brand-label">Social listening</div>
    <div class="page-title">Social Mention Analysis</div>
  </div>
  <div class="period-badge">YouTube &nbsp;·&nbsp; Q1 2026</div>
</div>

<!-- KPI band -->
<div class="kpi-row">
  <div class="kpi-card">
    <div class="kpi-value">{total_videos:,}</div>
    <div class="kpi-label">Videos Analyzed</div>
  </div>
  <div class="kpi-card">
    <div class="kpi-value">{total_comments:,}</div>
    <div class="kpi-label">Total Comments Analyzed</div>
  </div>
  <div class="kpi-card">
    <div class="kpi-value positive">{pos_pct}%</div>
    <div class="kpi-label">Positive Sentiment</div>
  </div>
</div>

<!-- Three-panel row -->
<div class="panels">

  <!-- Left: pie + big number -->
  <div class="panel-sentiment">
    <div class="panel-title">Overall Sentiment</div>
    <div class="chart-wrap">
      <canvas id="sentPie"></canvas>
    </div>
    <div class="legend">
      <div class="legend-item"><div class="legend-dot" style="background:#27ae60"></div> Positive</div>
      <div class="legend-item"><div class="legend-dot" style="background:#bdc3c7"></div> Neutral</div>
      <div class="legend-item"><div class="legend-dot" style="background:#e74c3c"></div> Negative</div>
    </div>
    <div class="big-pct">
      {pos_pct}%
      <div class="big-pct-label">of comments are Positive</div>
    </div>
  </div>

  <!-- Middle: comment highlights -->
  <div class="panel-comments">
    <div class="panel-title">Comment Highlights</div>
    {highlight_cards}
  </div>

  <!-- Right: key takeaways -->
  <div class="panel-takeaways">
    <div class="panel-title">Key Takeaways</div>
    {takeaway_items}
  </div>

</div>

<div class="footer">Generated {pd.Timestamp.now().strftime('%B %d, %Y')} &nbsp;|&nbsp; Renault Korea – Filante Influencer Campaign</div>

<script>
new Chart(document.getElementById('sentPie'), {{
  type: 'pie',
  data: {{
    labels: ['Positive', 'Neutral', 'Negative'],
    datasets: [{{
      data: [{pos}, {neu}, {neg}],
      backgroundColor: ['#27ae60', '#bdc3c7', '#e74c3c'],
      borderWidth: 2,
      borderColor: '#fff'
    }}]
  }},
  options: {{
    plugins: {{
      legend: {{ display: false }},
      tooltip: {{
        callbacks: {{
          label: function(ctx) {{
            const total = ctx.dataset.data.reduce((a,b)=>a+b,0);
            const pct = Math.round(ctx.parsed / total * 100);
            return ` ${{ctx.label}}: ${{ctx.parsed}} (${{pct}}%)`;
          }}
        }}
      }}
    }}
  }}
}});
</script>
</body>
</html>"""

    output_path.write_text(html, encoding="utf-8")
    print(f"[Report] HTML report saved → {output_path}")


def main():
    print("=" * 60)
    print("Renault Filante – Influencer Comment Analysis")
    print("=" * 60)

    # 1. Parse Excel
    videos = extract_videos_from_excel(EXCEL_PATH)

    # 2. Fetch comments from YouTube
    yt = build_youtube_client(GOOGLE_API_KEY)
    all_comments: list[dict] = []

    print(f"\n[YouTube] Fetching comments for {len(videos)} videos…")
    for i, video in enumerate(videos, 1):
        print(f"  ({i}/{len(videos)}) {video['channel']} – {video['title'][:60]}…")
        comments = fetch_comments(yt, video["video_id"])
        for c in comments:
            c["video_title"] = video["title"]
            c["channel"] = video["channel"]
        all_comments.extend(comments)
        time.sleep(0.15)

    if not all_comments:
        print("No comments collected. Check API key and quota.")
        return

    comments_df = pd.DataFrame(all_comments)
    print(f"\n[YouTube] Total comments collected: {len(comments_df)}")

    # 3. Classify sentiment
    print("\n[Gemini] Classifying comment sentiments…")
    gemini = init_gemini(GOOGLE_API_KEY)
    comments_df = classify_all_comments(gemini, comments_df)

    # 4. Save comments CSV (sorted by likes desc)
    csv_path = OUTPUT_DIR / "comments_by_likes.csv"
    comments_df.sort_values("likes", ascending=False).to_csv(csv_path, index=False, encoding="utf-8-sig")
    print(f"[Output] Comments CSV saved → {csv_path}")

    # 5. Generate key takeaways
    print("\n[Gemini] Generating key takeaways…")
    takeaways = generate_key_takeaways(gemini, comments_df)

    # 6. Build PPTX report
    report_path = OUTPUT_DIR / "report.pptx"
    build_pptx_report(videos, comments_df, takeaways, report_path)

    # 7. Print summary
    counts = comments_df["sentiment"].value_counts()
    total = len(comments_df)
    print("\n" + "=" * 60)
    print("SUMMARY")
    print("=" * 60)
    print(f"  Videos analysed : {len(videos)}")
    print(f"  Comments fetched: {total}")
    for sent in ("Positive", "Neutral", "Negative"):
        n = int(counts.get(sent, 0))
        print(f"  {sent:<12}: {n:>5}  ({n/total*100:.1f}%)")
    print(f"\n  Outputs in: {OUTPUT_DIR}")
    print("=" * 60)


if __name__ == "__main__":
    main()
